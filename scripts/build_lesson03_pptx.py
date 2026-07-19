#!/usr/bin/env python3
"""Build Lesson 03 core + appendix PPTX and regenerate declarative slide sources.

Usage:
    python3 scripts/build_lesson03_pptx.py
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK_DIR = ROOT / "decks" / "greenplum-query-tuning-theory"
SLIDES_DIR = DECK_DIR / "slides"
PPTX_CORE_ONLY = (
    ROOT / "lessons" / "lesson-03" / "artifacts" / "greenplum-query-tuning-core.pptx"
)
PPTX_FULL = (
    ROOT / "lessons" / "lesson-03" / "artifacts" / "greenplum-query-tuning-theory.pptx"
)
PPTX_APPENDIX = (
    ROOT
    / "lessons"
    / "lesson-03"
    / "artifacts"
    / "greenplum-query-tuning-appendix.pptx"
)

APPENDIX_DIVIDER = {
    "anchor": "appendix",
    "kicker": "Appendix",
    "title": "Appendix / deep reference",
    "subtitle": "Сюда ведут кнопки «→ Appendix» из core. Возврат — через портал «← Вернуться».",
    "type": "cards",
    "cards": [
        ["Зачем", "Deep-dive после core-лекции или homework.", "green"],
        ["Содержание", "ORCA/Legacy internals, stats formulas, TEMP ON COMMIT, storage.", "blue"],
        ["Вернуться", "Кнопка «→ Appendix» сначала открывает портал с «← Вернуться».", "amber"],
        ["Стенд", "SQL demos в labs/greenplum-625/examples/lesson03-*.sql", "green"],
    ],
    "jumps": [
        {"label": "→ Карта appendix", "anchor": "appendix-map"},
        {"label": "→ Формулы CE", "anchor": "appendix-ce-summary"},
        {"label": "→ Spill", "anchor": "appendix-spill"},
    ],
}

W, H = 12192000, 6858000
C = {
    "bg": RGBColor(0xF7, 0xF7, 0xF4),
    "panel": RGBColor(0xFF, 0xFF, 0xFF),
    "text": RGBColor(0x20, 0x21, 0x23),
    "muted": RGBColor(0x6E, 0x6E, 0x68),
    "green": RGBColor(0x10, 0xA3, 0x7F),
    "blue": RGBColor(0x2F, 0x6F, 0xED),
    "amber": RGBColor(0xB7, 0x79, 0x1F),
    "red": RGBColor(0xB4, 0x23, 0x18),
    "line": RGBColor(0xD8, 0xD7, 0xD0),
}

_SCRIPTS_DIR = Path(__file__).resolve().parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))
from lesson03_appendix_slide_specs import APPENDIX_SLIDES
from lesson03_core_slide_specs import CORE_SLIDES
from lesson03_front_matter import GLOSSARY_SLIDES
from lesson03_glossary_catalog import TERMS
from lesson03_pptx_richtext import (
    CODE_WRAP,
    looks_like_plan,
    looks_like_sql,
    normalize_code_lines,
    write_linked_textframe,
    write_plan_codeframe,
    write_sql_codeframe,
)


def inch(value: float) -> int:
    return int(value * 914400)


# Vertical budget (slide H = 7.5"). Google Slides inflates Calibri/Courier line boxes.
FOOTER_Y = inch(6.92)
CHIP_H = inch(0.34)
CHIP_GAP = inch(0.10)
# Default content floor when there are no term/jump chips (above footer nav).
CONTENT_BOTTOM = inch(6.20)
# Line-height fudge after Google font substitution (panels / soft-breaks).
GOOGLE_LINE_FACTOR = 1.72
# Code uses exact spcPts (~1.15× font); keep a small safety margin for fit tests.
CODE_LINE_FACTOR = 1.28
PANEL_WRAP_COLS = 78
CODE_WRAP_COLS = CODE_WRAP


def _add_rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _add_round(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = C["line"]
    shape.adjustments[0] = 0.08
    return shape


def _add_text(
    slide, x, y, w, h, text, *, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, font="Calibri"
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    write_linked_textframe(
        tf,
        text,
        size=size,
        bold=bold,
        color=color or C["text"],
        font=font,
        align=align,
    )
    return box


def _card(slide, x, y, w, h, title, body, color_key, *, body_size: int = 14):
    _add_round(slide, x, y, w, h, C["panel"])
    _add_rect(slide, x + inch(0.18), y + inch(0.16), inch(0.4), inch(0.045), C[color_key])
    title_size = 16 if len(body) > 100 else 17
    _add_text(
        slide,
        x + inch(0.2),
        y + inch(0.28),
        w - inch(0.36),
        inch(0.38),
        title,
        size=title_size,
        bold=True,
    )
    # Keep body inside the card; Google conversion expands line height.
    body_top = y + inch(0.68)
    body_h = h - inch(0.82)
    if body_h < inch(0.6):
        body_h = inch(0.6)
    # Auto-shrink body font for dense cards.
    approx_lines = max(1, (len(body) + 42) // 43)
    size = body_size
    avail = (body_h / 914400) * 72
    while size > 10 and approx_lines * size * 1.35 > avail:
        size -= 1
    _add_text(
        slide,
        x + inch(0.2),
        body_top,
        w - inch(0.36),
        body_h,
        body,
        size=size,
        color=C["muted"],
    )


def _nav_chip(slide, x, y, w, h, label: str, *, fill_key: str = "panel"):
    shape = _add_round(slide, x, y, w, h, C[fill_key] if fill_key != "panel" else C["panel"])
    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.auto_size = None
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = C["green"] if fill_key == "panel" else C["text"]
    tf.paragraphs[0].space_before = Pt(4)
    return shape


def _is_hub_slide(spec: dict) -> bool:
    """Shared hubs: TOC / teaching glossary / appendix entry (no personal portal)."""
    if spec.get("_is_portal") or spec.get("_is_portal_divider"):
        return True
    anchor = spec.get("anchor") or ""
    kicker = (spec.get("kicker") or "").lower()
    if anchor in {
        "toc",
        "toc-2",
        "glossary",
        "appendix",
        "appendix-map",
        "title",
        "nav-guide",
    }:
        return True
    if anchor.startswith("glossary") and not anchor.startswith("glossary-portal"):
        return True
    # Inline «теория → детали» pages are shared teaching hubs (like front glossary).
    if anchor.startswith("detail-"):
        return True
    if kicker.startswith("словарь ·") or kicker.startswith("оглавление"):
        return True
    if kicker.startswith("детали ·"):
        return True
    if "словарь рядом" in kicker:
        return True
    if kicker.startswith("appendix"):
        return True
    return False


def _set_click_tooltip(shape, tip: str) -> None:
    """OOXML hover hint: ``a:hlinkClick/@tooltip`` + ``cNvPr/@descr`` (alt text)."""
    if not tip:
        return
    text = tip[:900]
    c_nv_pr = shape._element.nvSpPr.cNvPr
    c_nv_pr.set("descr", text)
    hlink = c_nv_pr.find(qn("a:hlinkClick"))
    if hlink is not None:
        hlink.set("tooltip", text)


def _quick_glossary_cards() -> List[List[str]]:
    """Compact glossary for personal return-portals (Google-safe hard links)."""
    return [
        [
            "Selectivity",
            "sel ∈ (0,1]; rows ≈ N·sel. Ошибка оценки → плохой join/Motion.",
            "green",
        ],
        [
            "MCV / n_distinct",
            "MCV+freqs для частых значений; иначе sel ≈ 1/NDV. См. полный словарь.",
            "blue",
        ],
        [
            "Histogram",
            "Корзины равного числа строк; границы = разрезы. Для range (BETWEEN, >).",
            "amber",
        ],
        [
            "Motion / ORCA",
            "Redistribute · Broadcast · Gather. ORCA: optimizer=on.",
            "red",
        ],
    ]


def attach_nav_portals(
    specs: List[Dict[str, Any]],
) -> Tuple[List[Dict[str, Any]], Dict[int, int]]:
    """Append personal glossary/appendix portals with hardwired «← Вернуться».

    Google Slides breaks ``lastslideviewed`` (often acts like Next). Portals use
    ordinary slide links: origin → portal → origin.
    """
    base: List[Dict[str, Any]] = [dict(s) for s in specs]
    glossary_portals: List[Tuple[int, Dict[str, Any]]] = []
    jump_portals: List[Dict[str, Any]] = []
    has_glossary = any(s.get("anchor") == "glossary" for s in base)

    for index, spec in enumerate(base):
        if not _is_hub_slide(spec):
            origin_title = (spec.get("title") or "слайд")[:70]
            gloss_jumps = (
                [{"label": "→ Полный словарь 1/9", "anchor": "glossary"}]
                if has_glossary
                else []
            )
            glossary_portals.append(
                (
                    index,
                    {
                        "kicker": "Словарь",
                        "title": "Словарь (быстрый доступ)",
                        "subtitle": (
                            f"Открыто из: «{origin_title}». "
                            "«← Вернуться» вернёт на тот слайд."
                        ),
                        "type": "cards",
                        "cards": _quick_glossary_cards(),
                        "jumps": gloss_jumps,
                        "_return_to": index,
                        "_is_portal": True,
                    },
                )
            )

        new_jumps: List[Dict[str, str]] = []
        for jump_i, jump in enumerate(spec.get("jumps") or []):
            portal_anchor = f"nav-jump-{index}-{jump_i}"
            target_anchor = jump["anchor"]
            # Don't wrap portals that already point at other portals.
            if target_anchor.startswith("nav-jump-") or target_anchor.startswith(
                "glossary-portal"
            ):
                new_jumps.append(jump)
                continue
            origin_title = (spec.get("title") or "слайд")[:60]
            jump_portals.append(
                {
                    "anchor": portal_anchor,
                    "kicker": "Переход · возврат",
                    "title": jump.get("label", "→ Appendix").lstrip("→ ").strip(),
                    "subtitle": (
                        f"Из: «{origin_title}». Откройте раздел или вернитесь назад."
                    ),
                    "type": "cards",
                    "cards": [
                        [
                            "Открыть",
                            "Кнопка ниже ведёт в нужный раздел appendix/словаря.",
                            "green",
                        ],
                        [
                            "Вернуться",
                            "«← Вернуться» в футере — строго на слайд, с которого пришли.",
                            "blue",
                        ],
                        [
                            "Зачем портал",
                            "В Google Slides нет надёжного «last viewed» — только явная ссылка.",
                            "amber",
                        ],
                        [
                            "Меню",
                            "«☰ Меню» — в оглавление урока.",
                            "green",
                        ],
                    ],
                    "jumps": [
                        {
                            "label": jump.get("label") or f"→ {target_anchor}",
                            "anchor": target_anchor,
                        }
                    ],
                    "_return_to": index,
                    "_is_portal": True,
                }
            )
            new_jumps.append(
                {"label": jump.get("label") or f"→ {target_anchor}", "anchor": portal_anchor}
            )
        if new_jumps:
            base[index]["jumps"] = new_jumps

    if not glossary_portals and not jump_portals:
        return base, {}

    out = list(base)
    out.append(
        {
            "kicker": "Навигация",
            "title": "Служебные слайды возврата",
            "subtitle": "Не листать вручную — сюда ведут «Аа Словарь» и «→ Appendix».",
            "type": "cards",
            "cards": [
                [
                    "Словарь",
                    "У каждого слайда урока — свой портал с «← Вернуться».",
                    "green",
                ],
                [
                    "Appendix",
                    "Кнопка «→ Appendix» тоже через портал с возвратом.",
                    "blue",
                ],
                [
                    "Полный словарь",
                    "9 слайдов в начале презентации — из портала или оглавления.",
                    "amber",
                ],
                [
                    "Google Slides",
                    "Явные ссылки на слайды — возврат не зависит от истории просмотра.",
                    "red",
                ],
            ],
            "_is_portal_divider": True,
        }
    )

    origin_to_glossary_portal: Dict[int, int] = {}
    for origin_i, portal in glossary_portals:
        origin_to_glossary_portal[origin_i] = len(out)
        out.append(portal)
    for portal in jump_portals:
        out.append(portal)

    # Silence unused import warning path — GLOSSARY_SLIDES documents source of truth.
    _ = len(GLOSSARY_SLIDES)
    return out, origin_to_glossary_portal


def _slide_base(prs, spec, *, show_footer_label: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, C["bg"])
    _add_rect(slide, 0, 0, W, inch(0.02), C["line"])
    _add_rect(slide, 0, H - inch(0.02), W, inch(0.02), C["line"])
    # Compact header for dense content types so body stays above footer.
    dense = spec.get("type") in {"code", "panel", "image"}
    title_size = 24 if dense else (26 if spec.get("type") == "toc" else 28)
    title_h = inch(0.7) if dense else inch(0.95)
    subtitle_y = inch(1.45) if dense else (inch(1.7) if spec.get("type") == "toc" else inch(1.75))
    _add_text(
        slide,
        inch(0.55),
        inch(0.28),
        inch(8.5),
        inch(0.28),
        spec["kicker"].upper(),
        size=12,
        bold=True,
        color=C["green"],
    )
    _add_text(
        slide,
        inch(0.55),
        inch(0.58),
        inch(11.5),
        title_h,
        spec["title"],
        size=title_size,
        bold=True,
        font="Calibri",
    )
    if spec.get("subtitle"):
        _add_text(
            slide,
            inch(0.55),
            subtitle_y,
            inch(11.2),
            inch(0.45) if dense else inch(0.5),
            spec["subtitle"],
            size=14 if dense else 15,
            color=C["muted"],
        )
    if show_footer_label:
        _add_text(
            slide,
            inch(10.2),
            inch(7.05),
            inch(2.5),
            inch(0.3),
            "Урок 03",
            size=11,
            color=C["muted"],
            align=PP_ALIGN.RIGHT,
        )
    return slide


def _effective_code_lines(code: str, *, cols: int = CODE_WRAP_COLS) -> int:
    """Count visual lines after the same hard-wrap used in PPTX rendering."""
    if cols == CODE_WRAP:
        return max(1, len(normalize_code_lines(code, width=cols)))
    total = 0
    for line in code.split("\n"):
        total += max(1, (len(line) + cols - 1) // cols) if line else 1
    return max(1, total)


def _fit_font_size(
    lines: int,
    box_h_inch: float,
    *,
    candidates: tuple[int, ...],
    line_factor: float = GOOGLE_LINE_FACTOR,
) -> int:
    """Largest font that still fits ``lines`` into ``box_h_inch`` under Google inflation."""
    avail_pt = max(8.0, box_h_inch * 72.0)
    for size in candidates:
        if lines * size * line_factor <= avail_pt:
            return size
    return candidates[-1]


def _code_font_size(code: str, *, core: bool, box_h_inch: float = 3.15) -> int:
    """Pick a font that fits the code box after Google Slides font substitution."""
    lines = _effective_code_lines(code, cols=CODE_WRAP_COLS)
    candidates = (12, 11, 10, 9, 8, 7) if core else (11, 10, 9, 8, 7)
    return _fit_font_size(
        lines, box_h_inch, candidates=candidates, line_factor=CODE_LINE_FACTOR
    )


def _code_kind(spec: dict) -> str:
    explicit = spec.get("code_kind")
    if explicit in {"PLAN", "SQL"}:
        return explicit
    kicker = spec.get("kicker", "")
    if "Plan" in kicker or "Phases" in kicker or "EXPLAIN" in kicker:
        return "PLAN"
    code = spec.get("code", "")
    if looks_like_plan(code) and not looks_like_sql(code):
        return "PLAN"
    if looks_like_sql(code):
        return "SQL"
    if looks_like_plan(code):
        return "PLAN"
    return "SQL"


def _render_toc(slide, spec: dict) -> List[Tuple[Any, str, str]]:
    """Draw TOC entries; return (shape, target_anchor, tip) for click wiring."""
    entries = spec.get("entries") or []
    links: List[Tuple[Any, str, str]] = []
    card_w = inch(5.85)
    card_h = inch(0.72)
    gap_x = inch(0.3)
    gap_y = inch(0.12)
    origin_x = inch(0.55)
    origin_y = inch(2.35)
    for index, entry in enumerate(entries):
        col = index % 2
        row = index // 2
        x = origin_x + col * (card_w + gap_x)
        y = origin_y + row * (card_h + gap_y)
        color_key = entry.get("color", "green")
        shape = _add_round(slide, x, y, card_w, card_h, C["panel"])
        _add_rect(slide, x, y, inch(0.12), card_h, C[color_key])
        num = entry.get("num", "")
        title = entry.get("title", "")
        hint = entry.get("hint", "")
        box = slide.shapes.add_textbox(x + inch(0.28), y + inch(0.08), card_w - inch(0.4), card_h - inch(0.1))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.clear()
        r0 = p0.add_run()
        r0.text = f"{num}  {title}"
        r0.font.size = Pt(15)
        r0.font.bold = True
        r0.font.name = "Calibri"
        r0.font.color.rgb = C["text"]
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = hint
        r1.font.size = Pt(11)
        r1.font.name = "Calibri"
        r1.font.color.rgb = C["muted"]
        # Invisible hit target on the rounded card
        hit = shape
        links.append((hit, entry.get("anchor", "toc"), ""))
    return links


def _chip_rows(spec: dict) -> int:
    """How many chip rows sit between content and footer nav (0–2)."""
    return int(bool(spec.get("terms"))) + int(bool(spec.get("jumps")))


def _content_bottom_for(spec: dict) -> int:
    """Bottom edge of teaching content — always clear of term/jump chips + footer."""
    rows = _chip_rows(spec)
    if rows == 0:
        return CONTENT_BOTTOM
    # Each chip row + gap under the panel; keep clear air above FOOTER_Y.
    stack = rows * (CHIP_H + CHIP_GAP) + inch(0.14)
    return FOOTER_Y - stack


def _jump_row_y(spec: dict) -> int:
    """Y for appendix jump chips (closest chip row above footer)."""
    return FOOTER_Y - CHIP_H - CHIP_GAP


def _term_row_y(spec: dict) -> int:
    """Y for glossary term chips (above jumps when both present)."""
    if spec.get("jumps"):
        return _jump_row_y(spec) - CHIP_H - CHIP_GAP
    return _jump_row_y(spec)


def _prepare_textframe(tf) -> None:
    """Disable autofit and tighten insets — Google otherwise grows the box."""
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    try:
        body_pr = tf._txBody.bodyPr
        # EMUs: ~0.05" margins
        for attr, value in (
            ("lIns", "45720"),
            ("rIns", "45720"),
            ("tIns", "36576"),
            ("bIns", "36576"),
        ):
            body_pr.set(attr, value)
    except Exception:
        pass


def _render_terms(
    slide, terms: List[Dict[str, str]], *, y: int
) -> List[Tuple[Any, str, str]]:
    """Clickable term chips: hover tip + click → detail/glossary slide."""
    links: List[Tuple[Any, str, str]] = []
    if not terms:
        return links
    x = inch(0.55)
    max_right = inch(12.6)
    # Compact leading tag (no separate tall textbox that collides with panels).
    tag_chip = _nav_chip(slide, x, y, inch(0.95), CHIP_H, "Термины")
    for para in tag_chip.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = C["muted"]
            run.font.size = Pt(10)
    x += inch(1.05)
    for term in terms[:5]:
        label = term.get("label") or term.get("key") or "?"
        tip = term.get("tip") or ""
        anchor = term.get("anchor") or ""
        if not anchor:
            continue
        width = inch(min(2.45, max(1.1, 0.12 * len(label) + 0.5)))
        if x + width > max_right:
            break
        chip = _nav_chip(slide, x, y, width, CHIP_H, label)
        for para in chip.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = C["blue"]
        links.append((chip, anchor, tip))
        x += width + inch(0.08)
    return links


def _render_jumps(slide, jumps: List[Dict[str, str]], *, y: int) -> List[Tuple[Any, str, str]]:
    """Clickable «→ Appendix: …» chips just above the footer nav."""
    links: List[Tuple[Any, str, str]] = []
    if not jumps:
        return links
    x = inch(0.55)
    max_right = inch(12.6)
    for jump in jumps[:3]:
        label = jump.get("label") or f"→ {jump.get('anchor', 'appendix')}"
        width = inch(min(5.5, max(2.3, 0.11 * len(label) + 0.85)))
        if x + width > max_right:
            break
        chip = _nav_chip(slide, x, y, width, CHIP_H, label)
        links.append((chip, jump["anchor"], ""))
        x += width + inch(0.10)
    return links


def _render_panel(slide, spec: dict, *, core: bool, bottom: int) -> None:
    """Single teaching panel — long wrapped text that must stay above chip strips."""
    top = inch(2.02)
    panel_h = bottom - top
    if panel_h < inch(2.4):
        panel_h = inch(2.4)
        top = bottom - panel_h
    _add_round(slide, inch(0.55), top, inch(12.1), panel_h, C["panel"])
    body = spec.get("body") or ""
    lines = _effective_code_lines(body, cols=PANEL_WRAP_COLS)
    # Inner text box shorter than panel so Google reflow cannot paint over chips.
    pad_top = inch(0.14)
    pad_bottom = inch(0.18)
    text_h = panel_h - pad_top - pad_bottom
    box_h_inch = text_h / 914400
    candidates = (14, 13, 12, 11, 10, 9, 8) if core else (12, 11, 10, 9, 8)
    size = _fit_font_size(lines, box_h_inch, candidates=candidates)
    box = slide.shapes.add_textbox(inch(0.75), top + pad_top, inch(11.7), text_h)
    _prepare_textframe(box.text_frame)
    write_linked_textframe(box.text_frame, body, size=size, color=C["text"], font="Calibri")


def render_slide(prs, spec, *, core: bool = True):
    slide = _slide_base(prs, spec)
    toc_links: List[Tuple[Any, str, str]] = []
    bottom = _content_bottom_for(spec)
    if spec["type"] == "toc":
        toc_links = _render_toc(slide, spec)
    elif spec["type"] == "panel":
        _render_panel(slide, spec, core=core, bottom=bottom)
    elif spec["type"] == "code":
        label = _code_kind(spec)
        top = inch(2.02)
        panel_h = bottom - top
        label_h = inch(0.26)
        text_top = top + inch(0.34)
        text_h = panel_h - inch(0.48)
        text_h_inch = max(1.5, text_h / 914400)
        font_size = _code_font_size(spec["code"], core=core, box_h_inch=text_h_inch)
        _add_round(slide, inch(0.55), top, inch(12.1), panel_h, C["panel"])
        _add_text(
            slide,
            inch(0.8),
            top + inch(0.06),
            inch(2),
            label_h,
            label,
            size=11,
            bold=True,
            color=C["green"],
        )
        box = slide.shapes.add_textbox(inch(0.75), text_top, inch(11.7), text_h)
        tf = box.text_frame
        _prepare_textframe(tf)
        # Code writer sets word_wrap=False; keep autofit off.
        if label == "PLAN":
            write_plan_codeframe(tf, spec["code"], size=font_size)
        else:
            write_sql_codeframe(tf, spec["code"], size=font_size)
        tf.word_wrap = False
    elif spec["type"] == "image":
        image_path = ROOT / spec["image"]
        if not image_path.is_file():
            raise FileNotFoundError(f"Missing plan screenshot: {image_path}")
        left, top = inch(0.55), inch(2.05)
        max_w, max_h = inch(12.1), bottom - top
        slide.shapes.add_picture(str(image_path), left, top, width=max_w)
        pic = slide.shapes[-1]
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.height = max_h
            pic.width = int(pic.width * ratio)
            pic.left = int((W - pic.width) / 2)
    elif spec["type"] == "two":
        left, right = spec["left"], spec["right"]
        body = 13 if core else 12
        card_h = bottom - inch(2.2)
        _card(slide, inch(0.55), inch(2.2), inch(5.9), card_h, left[0], left[1], left[2], body_size=body)
        _card(slide, inch(6.75), inch(2.2), inch(5.9), card_h, right[0], right[1], right[2], body_size=body)
    elif spec["type"] == "flow":
        width = inch(2.05)
        gap = inch(0.25)
        y = inch(2.25)
        card_h = bottom - y
        for index, (title, body, color_key) in enumerate(spec["flow"]):
            x = inch(0.55) + index * (width + gap)
            _card(slide, x, y, width, card_h, title, body, color_key, body_size=12 if core else 11)
            if index < len(spec["flow"]) - 1:
                _add_rect(slide, x + width + inch(0.05), y + inch(1.5), gap - inch(0.1), inch(0.04), C["line"])
    else:
        cards = spec.get("cards") or []
        long_bodies = any(len(body) > 80 for _, body, _ in cards)
        # 2×2 grid must end above footer / jump / terms strip.
        top_y = inch(2.15)
        gap_y = inch(0.12)
        card_h = min(inch(1.95) if long_bodies else inch(1.75), (bottom - top_y - gap_y) // 2)
        body_size = 12 if long_bodies else (13 if core else 12)
        positions = [
            (inch(0.55), top_y),
            (inch(6.75), top_y),
            (inch(0.55), top_y + card_h + gap_y),
            (inch(6.75), top_y + card_h + gap_y),
        ]
        for index, (title, body, color_key) in enumerate(cards):
            x, y = positions[index]
            _card(slide, x, y, inch(5.9), card_h, title, body, color_key, body_size=body_size)
    if spec.get("terms"):
        toc_links.extend(_render_terms(slide, spec.get("terms") or [], y=_term_row_y(spec)))
    if spec.get("jumps"):
        toc_links.extend(_render_jumps(slide, spec.get("jumps") or [], y=_jump_row_y(spec)))
    return slide, toc_links


def _anchor_index(specs: List[Dict[str, Any]]) -> Dict[str, int]:
    out: Dict[str, int] = {}
    for index, spec in enumerate(specs):
        anchor = spec.get("anchor")
        if anchor and anchor not in out:
            out[anchor] = index
    return out


def _wire_navigation(
    prs: Presentation,
    specs: List[Dict[str, Any]],
    toc_link_map: List[List[Tuple[Any, str, str]]],
    origin_to_glossary_portal: Optional[Dict[int, int]] = None,
) -> None:
    """Footer chips: Menu / Словарь / hardwired «← Вернуться» on portals."""
    origin_to_glossary_portal = origin_to_glossary_portal or {}
    anchors = _anchor_index(specs)
    toc_idx = anchors.get("toc", 0)
    glossary_idx = anchors.get("glossary", toc_idx)
    slides = list(prs.slides)
    toc_slide = slides[toc_idx]
    glossary_slide = slides[glossary_idx]

    missing = sorted(
        {
            anchor
            for slide_links in toc_link_map
            for _shape, anchor, _tip in slide_links
            if anchor not in anchors
        }
    )
    # Full deck must resolve every jump; core-only omits appendix by design.
    has_appendix = "appendix" in anchors or "appendix-map" in anchors
    if missing and has_appendix:
        raise RuntimeError(f"Unknown slide anchors in jumps/TOC: {missing}")
    # Core-only still must resolve term chips → detail-* / glossary.
    if missing and not has_appendix:
        detail_missing = [a for a in missing if a.startswith("detail-") or a.startswith("glossary")]
        if detail_missing:
            raise RuntimeError(f"Unknown glossary/detail anchors: {detail_missing}")

    for index, slide in enumerate(slides):
        spec = specs[index]
        menu = _nav_chip(slide, inch(0.45), FOOTER_Y, inch(1.45), inch(0.38), "☰ Меню")
        menu.click_action.target_slide = toc_slide

        if spec.get("_is_portal"):
            # Personal portal: Словарь → full glossary; Вернуться → exact origin.
            gloss = _nav_chip(slide, inch(2.0), FOOTER_Y, inch(1.55), inch(0.38), "Аа Словарь")
            gloss.click_action.target_slide = glossary_slide
            back = _nav_chip(slide, inch(3.65), FOOTER_Y, inch(1.95), inch(0.38), "← Вернуться")
            return_to = spec.get("_return_to")
            if isinstance(return_to, int) and 0 <= return_to < len(slides):
                back.click_action.target_slide = slides[return_to]
            continue

        if spec.get("_is_portal_divider"):
            continue

        gloss = _nav_chip(slide, inch(2.0), FOOTER_Y, inch(1.55), inch(0.38), "Аа Словарь")
        portal_i = origin_to_glossary_portal.get(index)
        if portal_i is not None:
            gloss.click_action.target_slide = slides[portal_i]
        else:
            # Shared hubs (TOC / teaching glossary): open the real glossary.
            gloss.click_action.target_slide = glossary_slide

        # Shared glossary pages: explicit back to TOC (no lastslideviewed).
        if (spec.get("anchor") or "").startswith("glossary") or (
            (spec.get("kicker") or "").startswith("Словарь ·")
        ):
            back = _nav_chip(slide, inch(3.65), FOOTER_Y, inch(1.95), inch(0.38), "← К оглавлению")
            back.click_action.target_slide = toc_slide

        # Inline detail hubs: jump back to the matching front glossary page.
        anchor = spec.get("anchor") or ""
        if anchor.startswith("detail-"):
            key = anchor[len("detail-") :]
            front = TERMS.get(key, {}).get("front_anchor", "glossary")
            front_i = anchors.get(front, glossary_idx)
            back = _nav_chip(slide, inch(3.65), FOOTER_Y, inch(2.15), inch(0.38), "← В словарь")
            back.click_action.target_slide = slides[front_i]

        # Explicit journey chips — do not dump live viewers into the full glossary.
        if anchor == "nav-guide":
            start = slides[anchors.get("stage-problem", min(index + 1, len(slides) - 1))]
            go = _nav_chip(
                slide, inch(3.65), FOOTER_Y, inch(2.35), inch(0.38), "Core 60 →"
            )
            go.click_action.target_slide = start
        elif spec.get("type") == "toc":
            # toc → toc-2; toc-2 → stage-problem (skip front glossary for live).
            nxt_anchor = "toc-2" if anchor == "toc" else "stage-problem"
            nxt = slides[anchors.get(nxt_anchor, min(index + 1, len(slides) - 1))]
            go = _nav_chip(slide, inch(3.65), FOOTER_Y, inch(1.55), inch(0.38), "Дальше →")
            go.click_action.target_slide = nxt

    # Wire TOC / term chips / appendix jumps (+ hover tooltips on terms).
    for slide_links in toc_link_map:
        for shape, anchor, tip in slide_links:
            target_i = anchors.get(anchor)
            if target_i is None:
                continue
            shape.click_action.target_slide = slides[target_i]
            _set_click_tooltip(shape, tip)


def write_shared_mjs() -> None:
    shared = (ROOT / "decks" / "greenplum-partitioning-theory" / "slides" / "shared.mjs").read_text(
        encoding="utf-8"
    )
    shared = shared.replace("Урок 02", "Урок 03")
    shared = shared.replace(
        '  if (spec.type === "code") {\n    codeBlock(ctx, slide, 84, 260, 1060, 310, spec.code);',
        '  if (spec.type === "image") {\n'
        '    codeBlock(ctx, slide, 84, 260, 1060, 310,\n'
        '      "Скрин EXPLAIN (см. PPTX / artifacts):\\n" + (spec.image || ""));\n'
        '  } else if (spec.type === "code") {\n'
        '    codeBlock(ctx, slide, 84, 260, 1060, 310, spec.code);',
    )
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    (SLIDES_DIR / "shared.mjs").write_text(shared, encoding="utf-8")


def write_content_mjs(slides, path: Path) -> None:
    def color_expr(key: str) -> str:
        return f"C.{key}"

    chunks = ["import { C } from \"./shared.mjs\";", "", "export const slides = ["]
    for spec in slides:
        chunks.append("  {")
        chunks.append(f'    kicker: {json.dumps(spec["kicker"], ensure_ascii=False)},')
        chunks.append(f'    title: {json.dumps(spec["title"], ensure_ascii=False)},')
        chunks.append(f'    subtitle: {json.dumps(spec.get("subtitle", ""), ensure_ascii=False)},')
        chunks.append(f'    type: {json.dumps(spec["type"], ensure_ascii=False)},')
        if spec.get("anchor"):
            chunks.append(f'    anchor: {json.dumps(spec["anchor"], ensure_ascii=False)},')
        if spec["type"] == "code":
            chunks.append(f'    code: {json.dumps(spec["code"], ensure_ascii=False)},')
        elif spec["type"] == "panel":
            chunks.append(f'    body: {json.dumps(spec.get("body", ""), ensure_ascii=False)},')
        elif spec["type"] == "image":
            chunks.append(f'    image: {json.dumps(spec["image"], ensure_ascii=False)},')
        elif spec["type"] == "toc":
            chunks.append(f'    entries: {json.dumps(spec.get("entries", []), ensure_ascii=False)},')
        elif spec["type"] == "two":
            left = spec["left"]
            right = spec["right"]
            chunks.append(
                f'    left: [{json.dumps(left[0], ensure_ascii=False)}, {json.dumps(left[1], ensure_ascii=False)}, {color_expr(left[2])}],'
            )
            chunks.append(
                f'    right: [{json.dumps(right[0], ensure_ascii=False)}, {json.dumps(right[1], ensure_ascii=False)}, {color_expr(right[2])}],'
            )
        elif spec["type"] == "flow":
            chunks.append("    flow: [")
            for title, body, color_key in spec["flow"]:
                chunks.append(
                    f'      [{json.dumps(title, ensure_ascii=False)}, {json.dumps(body, ensure_ascii=False)}, {color_expr(color_key)}],'
                )
            chunks.append("    ],")
        else:
            chunks.append("    cards: [")
            for title, body, color_key in spec.get("cards") or []:
                chunks.append(
                    f'      [{json.dumps(title, ensure_ascii=False)}, {json.dumps(body, ensure_ascii=False)}, {color_expr(color_key)}],'
                )
            chunks.append("    ],")
        if spec.get("jumps"):
            chunks.append(f'    jumps: {json.dumps(spec["jumps"], ensure_ascii=False)},')
        if spec.get("terms"):
            chunks.append(f'    terms: {json.dumps(spec["terms"], ensure_ascii=False)},')
        chunks.append("  },")
    chunks.append("];")
    chunks.append("")
    path.write_text("\n".join(chunks), encoding="utf-8")


def write_slide_modules(count: int) -> None:
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    for stale in SLIDES_DIR.glob("slide-*.mjs"):
        stale.unlink()
    for index in range(1, count + 1):
        name = f"slide-{index:02d}.mjs"
        fn = f"slide{index:02d}"
        body = (
            'import { slides } from "./content.mjs";\n'
            'import { renderContentSlide } from "./shared.mjs";\n\n'
            f"export async function {fn}(presentation, ctx) {{\n"
            f"  return renderContentSlide(presentation, ctx, slides[{index - 1}]);\n"
            "}\n"
        )
        (SLIDES_DIR / name).write_text(body, encoding="utf-8")


def build_pptx(
    slides: List[Dict[str, Any]],
    path: Path,
    *,
    core: bool,
    origin_to_glossary_portal: Optional[Dict[int, int]] = None,
) -> None:
    prs = Presentation()
    prs.slide_width = Emu(W)
    prs.slide_height = Emu(H)
    toc_link_map: List[List[Tuple[Any, str]]] = []
    for spec in slides:
        _slide, links = render_slide(prs, spec, core=core)
        toc_link_map.append(links)
    _wire_navigation(prs, slides, toc_link_map, origin_to_glossary_portal)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def main() -> None:
    write_shared_mjs()
    core_slides, core_portals = attach_nav_portals(list(CORE_SLIDES))
    appendix_slides, appendix_portals = attach_nav_portals(list(APPENDIX_SLIDES))
    full_base = list(CORE_SLIDES) + [APPENDIX_DIVIDER] + list(APPENDIX_SLIDES)
    full_slides, full_portals = attach_nav_portals(full_base)

    write_content_mjs(full_slides, SLIDES_DIR / "content.mjs")
    write_slide_modules(len(full_slides))
    build_pptx(
        core_slides,
        PPTX_CORE_ONLY,
        core=True,
        origin_to_glossary_portal=core_portals,
    )
    build_pptx(
        appendix_slides,
        PPTX_APPENDIX,
        core=False,
        origin_to_glossary_portal=appendix_portals,
    )
    build_pptx(
        full_slides,
        PPTX_FULL,
        core=True,
        origin_to_glossary_portal=full_portals,
    )
    write_content_mjs(appendix_slides, SLIDES_DIR / "content-appendix.mjs")
    print(f"Wrote {len(core_slides)} core-only slides to {PPTX_CORE_ONLY}")
    print(f"Wrote {len(appendix_slides)} appendix slides to {PPTX_APPENDIX}")
    print(
        f"Wrote {len(full_slides)} full slides (core+divider+appendix+portals) to {PPTX_FULL}"
    )
    print(
        f"Glossary portals: core={len(core_portals)} full={len(full_portals)} "
        f"(hardwired ← Вернуться; Google-safe)"
    )
    print(f"Wrote sources under {SLIDES_DIR}")


if __name__ == "__main__":
    main()
