"""Rich text helpers for Lesson 03 PPTX: clickable URLs + SQL/PLAN highlighting.

Google Slides note
------------------
Soft line breaks (``a:br``) + many colored runs make SQL look «crooked»
(uneven gaps, collapsed indents). Code frames therefore use:

- **one paragraph per visual line**;
- **exact** ``spcPts`` line spacing (not percentage);
- leading spaces → NBSP so Google keeps indentation;
- optional hard-wrap at ``CODE_WRAP`` so ``word_wrap=False`` is safe.
"""

from __future__ import annotations

import re
from typing import List, Optional, Tuple

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
from pygments import lex
from pygments.lexers import SqlLexer
from pygments.token import Comment, Keyword, Literal, Name, Number, Operator, Punctuation, String, Text, Token

URL_RE = re.compile(r"https?://[^\s<>\"')\]]+")

SQL_COLORS = {
    "keyword": RGBColor(0x2F, 0x6F, 0xED),
    "string": RGBColor(0x10, 0xA3, 0x7F),
    "number": RGBColor(0xB7, 0x79, 0x1F),
    "comment": RGBColor(0x6E, 0x6E, 0x68),
    "name": RGBColor(0x20, 0x21, 0x23),
    "op": RGBColor(0x20, 0x21, 0x23),
    "default": RGBColor(0x20, 0x21, 0x23),
    "link": RGBColor(0x2F, 0x6F, 0xED),
}

PLAN_KEYWORDS = sorted(
    {
        "Nested Loop",
        "Hash Join",
        "Merge Join",
        "Hash Semi Join",
        "Hash Anti Join",
        "Hash Left Anti Semi (Not-In) Join",
        "Broadcast Motion",
        "Redistribute Motion",
        "Gather Motion",
        "Seq Scan",
        "Index Scan",
        "Bitmap Heap Scan",
        "HashAggregate",
        "GroupAggregate",
        "WindowAgg",
        "Materialize",
        "Values Scan",
        "Sort",
        "Result",
        "Partition By",
        "Order By",
        "Hash Key",
        "Hash Cond",
        "One-Time Filter",
        "Optimizer",
        "GPORCA",
        "Postgres",
    },
    key=len,
    reverse=True,
)
_PLAN_RE = re.compile("|".join(re.escape(n) for n in PLAN_KEYWORDS))

# Courier New converts more reliably in Google Slides than Consolas.
CODE_FONT = "Courier New"
CODE_WRAP = 68


def _style_run(run, *, size: int, bold: bool, color: RGBColor, font: str) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    # Force all script slots — Google otherwise substitutes a proportional font mid-line.
    try:
        rPr = run._r.get_or_add_rPr()
        for tag in ("latin", "ea", "cs"):
            el = rPr.find(qn(f"a:{tag}"))
            if el is None:
                el = etree.SubElement(rPr, qn(f"a:{tag}"))
            el.set("typeface", font)
    except Exception:
        pass


def _split_urls(text: str) -> List[Tuple[str, Optional[str]]]:
    parts: List[Tuple[str, Optional[str]]] = []
    pos = 0
    for match in URL_RE.finditer(text):
        if match.start() > pos:
            parts.append((text[pos : match.start()], None))
        raw = match.group(0)
        url = raw.rstrip(".,;")
        parts.append((url, url))
        trailing = raw[len(url) :]
        if trailing:
            parts.append((trailing, None))
        pos = match.end()
    if pos < len(text):
        parts.append((text[pos:], None))
    return parts or [(text, None)]


def _preserve_indent(text: str) -> str:
    """Keep leading whitespace in Google Slides (regular spaces often collapse)."""
    i = 0
    while i < len(text) and text[i] in " \t":
        i += 1
    if i == 0:
        return text
    lead = text[:i].expandtabs(2).replace(" ", "\u00A0")
    return lead + text[i:]


def _hard_wrap_line(line: str, width: int = CODE_WRAP) -> List[str]:
    """Wrap a long code line at spaces; continuation keeps +2 indent."""
    if len(line) <= width:
        return [line]
    indent = 0
    while indent < len(line) and line[indent] in " \t":
        indent += 1
    body = line[indent:]
    prefix = line[:indent].expandtabs(2)
    cont = prefix + "  "
    out: List[str] = []
    while body:
        if len(prefix) + len(body) <= width:
            out.append(prefix + body)
            break
        budget = max(8, width - len(prefix))
        cut = body.rfind(" ", 0, budget + 1)
        if cut < 4:
            cut = budget
        out.append(prefix + body[:cut].rstrip())
        body = body[cut:].lstrip()
        prefix = cont
    return out or [line]


def normalize_code_lines(code: str, *, width: int = CODE_WRAP) -> List[str]:
    """Split code into visual lines ready for one-paragraph-per-line rendering."""
    visual: List[str] = []
    for raw in (code or "").split("\n"):
        visual.extend(_hard_wrap_line(raw, width))
    return visual or [""]


def _set_exact_paragraph_spacing(paragraph, *, size_pt: int) -> None:
    """Lock spacing in absolute points — Google ignores % spacing unpredictably."""
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    p_pr = paragraph._p.get_or_add_pPr()
    for tag in ("spcBef", "spcAft", "lnSpc"):
        node = p_pr.find(qn(f"a:{tag}"))
        if node is not None:
            p_pr.remove(node)
    spc_bef = etree.SubElement(p_pr, qn("a:spcBef"))
    etree.SubElement(spc_bef, qn("a:spcPts")).set("val", "0")
    spc_aft = etree.SubElement(p_pr, qn("a:spcAft"))
    etree.SubElement(spc_aft, qn("a:spcPts")).set("val", "0")
    ln_spc = etree.SubElement(p_pr, qn("a:lnSpc"))
    # Slightly taller than font — enough for glyphs, not enough for Google inflation.
    etree.SubElement(ln_spc, qn("a:spcPts")).set("val", str(int(size_pt * 100 * 1.15)))


def _tighten_paragraph(paragraph) -> None:
    """Zero spacing for prose (soft-break path)."""
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    try:
        paragraph.line_spacing = 1.0
    except (AttributeError, TypeError, ValueError):
        pass


def _add_line_break(paragraph) -> None:
    """Soft line break inside one paragraph (prose only)."""
    etree.SubElement(paragraph._p, qn("a:br"))


def _merge_same_style(
    pieces: List[Tuple[str, RGBColor, Optional[str], bool]],
) -> List[Tuple[str, RGBColor, Optional[str], bool]]:
    """Collapse adjacent runs with identical style — fewer Google reflow glitches."""
    if not pieces:
        return pieces
    out = [pieces[0]]
    for value, color, url, bold in pieces[1:]:
        prev_v, prev_c, prev_u, prev_b = out[-1]
        if color == prev_c and url == prev_u and bold == prev_b and url is None:
            out[-1] = (prev_v + value, prev_c, prev_u, prev_b)
        else:
            out.append((value, color, url, bold))
    return out


def write_linked_textframe(
    text_frame,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    font: str = "Calibri",
    align=PP_ALIGN.LEFT,
) -> None:
    """Plain multi-line text with clickable http(s) URLs (single paragraph + breaks)."""
    text_frame.word_wrap = True
    base = color or SQL_COLORS["default"]
    lines = text.split("\n") if text is not None else [""]
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.clear()
    _tighten_paragraph(paragraph)
    for index, line in enumerate(lines):
        if index > 0:
            _add_line_break(paragraph)
        segments = _split_urls(line)
        wrote = False
        for fragment, url in segments:
            if fragment == "" and wrote:
                continue
            run = paragraph.add_run()
            run.text = fragment if fragment != "" else " "
            wrote = True
            _style_run(
                run,
                size=size,
                bold=bold,
                color=SQL_COLORS["link"] if url else base,
                font=font,
            )
            if url:
                run.hyperlink.address = url
        if not wrote:
            run = paragraph.add_run()
            run.text = " "
            _style_run(run, size=size, bold=bold, color=base, font=font)


def _token_color(ttype) -> RGBColor:
    if ttype in Comment or ttype in Token.Comment:
        return SQL_COLORS["comment"]
    if ttype in Keyword or ttype in Token.Keyword:
        return SQL_COLORS["keyword"]
    if ttype in String or ttype in Literal.String:
        return SQL_COLORS["string"]
    if ttype in Number or ttype in Literal.Number:
        return SQL_COLORS["number"]
    if ttype in Operator or ttype in Punctuation:
        return SQL_COLORS["op"]
    if ttype in Name or ttype in Text:
        return SQL_COLORS["name"]
    return SQL_COLORS["default"]


def _sql_line_pieces(line: str) -> List[Tuple[str, RGBColor, Optional[str], bool]]:
    pieces: List[Tuple[str, RGBColor, Optional[str], bool]] = []
    for fragment, url in _split_urls(line):
        if url:
            pieces.append((fragment, SQL_COLORS["link"], url, False))
            continue
        for ttype, value in lex(fragment, SqlLexer()):
            if value == "":
                continue
            for part in value.split("\n"):
                if part == "":
                    continue
                pieces.append((part, _token_color(ttype), None, False))
    return _merge_same_style(pieces) or [(" ", SQL_COLORS["default"], None, False)]


def _nbsp_leading(pieces: List[Tuple[str, RGBColor, Optional[str], bool]]) -> List[
    Tuple[str, RGBColor, Optional[str], bool]
]:
    """Convert leading spaces/tabs across the first runs into NBSP."""
    if not pieces:
        return [("\u00A0", SQL_COLORS["default"], None, False)]
    out: List[Tuple[str, RGBColor, Optional[str], bool]] = []
    leading = True
    for value, color, url, bold in pieces:
        if leading:
            i = 0
            while i < len(value) and value[i] in " \t":
                i += 1
            ws = value[:i].expandtabs(2).replace(" ", "\u00A0")
            rest = value[i:]
            if ws:
                out.append((ws, SQL_COLORS["default"], None, False))
            if rest or url is not None:
                if rest or url:
                    out.append((rest if rest else value, color, url, bold))
                leading = False
            # else: still in pure-whitespace runs
        else:
            out.append((value, color, url, bold))
    return out or [("\u00A0", SQL_COLORS["default"], None, False)]


def _write_code_lines(
    text_frame,
    lines: List[str],
    *,
    size: int,
    line_pieces,
) -> None:
    """One paragraph per line + exact spacing — stable in Google Slides."""
    text_frame.word_wrap = False
    first = True
    for line in lines:
        if first:
            paragraph = text_frame.paragraphs[0]
            paragraph.clear()
            first = False
        else:
            paragraph = text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        _set_exact_paragraph_spacing(paragraph, size_pt=size)
        pieces = _nbsp_leading(line_pieces(line if line != "" else " "))
        for value, color, url, bold in pieces:
            run = paragraph.add_run()
            run.text = value if value != "" else "\u00A0"
            _style_run(run, size=size, bold=bold, color=color, font=CODE_FONT)
            if url:
                run.hyperlink.address = url


def write_sql_codeframe(text_frame, code: str, *, size: int = 14) -> None:
    """Render SQL with Pygments colors — Google-stable paragraph-per-line layout."""
    lines = normalize_code_lines(code, width=CODE_WRAP)
    _write_code_lines(text_frame, lines, size=size, line_pieces=_sql_line_pieces)


def write_plan_codeframe(text_frame, code: str, *, size: int = 14) -> None:
    """Render EXPLAIN text — Google-stable paragraph-per-line layout."""

    def plan_pieces(line: str) -> List[Tuple[str, RGBColor, Optional[str], bool]]:
        pieces: List[Tuple[str, RGBColor, Optional[str], bool]] = []
        for frag, url in _split_urls(line):
            if url:
                pieces.append((frag, SQL_COLORS["link"], url, True))
                continue
            cursor = 0
            for match in _PLAN_RE.finditer(frag):
                if match.start() > cursor:
                    pieces.append(
                        (frag[cursor : match.start()], SQL_COLORS["default"], None, False)
                    )
                pieces.append((match.group(0), SQL_COLORS["keyword"], None, True))
                cursor = match.end()
            if cursor < len(frag):
                pieces.append((frag[cursor:], SQL_COLORS["default"], None, False))
        return _merge_same_style(pieces) or [(" ", SQL_COLORS["default"], None, False)]

    lines = normalize_code_lines(code, width=CODE_WRAP)
    _write_code_lines(text_frame, lines, size=size, line_pieces=plan_pieces)


def looks_like_plan(code: str) -> bool:
    markers = ("Motion", "Hash Join", "Nested Loop", "WindowAgg", "Seq Scan", "->")
    return any(m in code for m in markers) and "CREATE TABLE" not in code.upper()


def looks_like_sql(code: str) -> bool:
    upper = code.upper()
    return any(
        token in upper
        for token in (
            "SELECT ",
            "SELECT\n",
            "CREATE ",
            "WITH ",
            "INSERT ",
            "EXPLAIN",
            "DISTRIBUTED",
            "ALTER TABLE",
            "ANALYZE ",
            "SET ",
            "DROP ",
            "JOIN ",
            "WHERE ",
            "VALUES",
            "PARTITION BY",
            "NOT IN",
            "NOT EXISTS",
        )
    )
